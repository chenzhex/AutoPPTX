import pytest
from pptx import Presentation

from autopptx.Text.style import (
    set_paragraph_style,
    extract_paragraph_style,
    set_textbox_style,
    extract_textbox_style,
    transfer_text_style,
)
from autopptx.Type.find import find_placeholders

TEST_PPTX = "data/output_demo.pptx"


@pytest.fixture
def prs():
    """Load and return the test PowerPoint presentation object."""
    return Presentation(TEST_PPTX)


@pytest.fixture
def src_dst_shapes(prs):
    """
    Extract source and destination bodytext shapes from slide 0 and slide 1.
    Used for testing style transfer between placeholders.
    """
    src = find_placeholders(prs.slides[0], "bodytext")[0]
    dst = find_placeholders(prs.slides[1], "bodytext")[0]
    return src, dst


def test_set_paragraph_style_basic(src_dst_shapes):
    """
    Test basic font and alignment styling on a single paragraph.
    Verifies that each attribute is set correctly.
    """
    src_shape, _ = src_dst_shapes
    p = src_shape.text_frame.paragraphs[0]
    set_paragraph_style(
        p,
        font_name="Arial",
        font_size=20,
        bold=True,
        italic=True,
        font_color=(255, 0, 0),
        align="center",
    )
    style = extract_paragraph_style(p)

    assert style["font_name"] == "Arial"
    assert style["font_size"] == 20
    assert style["bold"] is True
    assert style["italic"] is True
    assert style["font_color"] == (255, 0, 0)
    assert style["align"] == "CENTER"


@pytest.mark.parametrize("align_value", ["left", "center", "right", "justify", "distribute"])
def test_set_paragraph_style_alignment(src_dst_shapes, align_value):
    """
    Parametrized test to verify paragraph alignment is applied correctly.
    Covers all valid alignment options.
    """
    src_shape, _ = src_dst_shapes
    p = src_shape.text_frame.paragraphs[0]
    set_paragraph_style(p, align=align_value)
    style = extract_paragraph_style(p)
    assert style["align"].lower() == align_value


def test_set_textbox_style_applies_all(src_dst_shapes):
    """
    Test whether textbox-level styling is applied to all paragraphs.
    Verifies that font name, size, and color are uniform.
    """
    src_shape, _ = src_dst_shapes
    set_textbox_style(src_shape, font_name="Calibri", font_size=22, font_color=(0, 128, 255))
    styles = extract_textbox_style(src_shape)
    assert all(s["font_name"] == "Calibri" for s in styles)
    assert all(s["font_size"] == 22 for s in styles)
    assert all(s["font_color"] == (0, 128, 255) for s in styles)


@pytest.mark.parametrize("mode", ["single", "full"])
def test_transfer_text_style(src_dst_shapes, mode):
    """
    Test transfer of text styles from source to destination shape.
    Verifies both 'single' and 'full' modes apply the expected styles.
    """
    src_shape, dst_shape = src_dst_shapes
    set_textbox_style(src_shape, font_name="Verdana", font_size=18, font_color=(10, 100, 200))
    transfer_text_style(src_shape, dst_shape, mode=mode)

    dst_styles = extract_textbox_style(dst_shape)
    assert all(style["font_name"] == "Verdana" for style in dst_styles)
    assert all(style["font_size"] == 18 for style in dst_styles)
    assert all(style["font_color"] == (10, 100, 200) for style in dst_styles)


def test_transfer_text_style_invalid_mode(src_dst_shapes):
    """
    Test that an unsupported mode in transfer_text_style raises ValueError.
    """
    src_shape, dst_shape = src_dst_shapes
    with pytest.raises(ValueError):
        transfer_text_style(src_shape, dst_shape, mode="unsupported")


def test_extract_paragraph_style_invalid():
    """
    Test that extract_paragraph_style raises an error for invalid input type.
    """
    with pytest.raises(ValueError):
        extract_paragraph_style("not a paragraph")


def test_set_paragraph_style_invalid():
    """
    Test that set_paragraph_style raises an error for non-paragraph input.
    """
    with pytest.raises(ValueError):
        set_paragraph_style("not a paragraph")


def test_set_paragraph_style_invalid_align(src_dst_shapes):
    """
    Test that setting an unknown alignment string raises ValueError.
    """
    src_shape, _ = src_dst_shapes
    p = src_shape.text_frame.paragraphs[0]
    with pytest.raises(ValueError):
        set_paragraph_style(p, align="unknown")