import pytest
from pptx import Presentation

from autopptx.Table.style import (
    set_table_style,
    extract_table_style,
    transfer_table_style,
    extract_cell_style
)

TEST_PPTX = "data/output_demo.pptx"


@pytest.fixture
def prs():
    return Presentation(TEST_PPTX)


@pytest.fixture
def table_pair(prs):
    """Return two tables from the 6th slide: source and destination."""
    slide = prs.slides[5]
    return slide.shapes[2], slide.shapes[3]


@pytest.fixture
def styled_table(prs):
    """Return a table that has been styled."""
    slide = prs.slides[5]
    table_shape = slide.shapes[2]
    set_table_style(
        table_shape,
        font_name="Arial",
        font_size=16,
        bold=True,
        italic=False,
        font_color=(255, 0, 0),
        align="center",
        bg_color=(200, 200, 200),
    )
    return table_shape


def test_set_table_style_applies_to_all_cells(styled_table):
    """Test that table style is correctly applied to all cells."""
    table = styled_table.table
    for row in table.rows:
        for cell in row.cells:
            style = extract_cell_style(cell)
            assert style["font_name"] == "Arial"
            assert style["font_size"] == 16
            assert style["bold"] is True
            assert style["italic"] is False
            assert style["font_color"] == (255, 0, 0)
            assert style["align"].lower() == "center"
            assert style["bg_color"] == (200, 200, 200)


def test_extract_table_style_returns_matrix(styled_table):
    """Test that extract_table_style returns a 2D list of styles."""
    table = styled_table.table
    styles = extract_table_style(styled_table)
    assert len(styles) == len(table.rows)
    for row in styles:
        assert len(row) == len(table.columns)
        for cell_style in row:
            assert isinstance(cell_style, dict)


def test_transfer_table_style_full(table_pair):
    """Test that transfer_table_style with mode='full' works."""
    src_shape, dst_shape = table_pair
    transfer_table_style(src_shape, dst_shape, mode="full")
    src_styles = extract_table_style(src_shape)
    dst_styles = extract_table_style(dst_shape)

    for i in range(min(len(src_styles), len(dst_styles))):
        for j in range(min(len(src_styles[i]), len(dst_styles[i]))):
            assert src_styles[i][j] == dst_styles[i][j]


def test_transfer_table_style_single(table_pair):
    """Test that transfer_table_style with mode='single' applies (0,0) style to all."""
    src_shape, dst_shape = table_pair
    transfer_table_style(src_shape, dst_shape, mode="single")
    style_00 = extract_cell_style(src_shape.table.cell(0, 0))
    dst_table = dst_shape.table

    for row in dst_table.rows:
        for cell in row.cells:
            style = extract_cell_style(cell)
            for k in style_00:
                assert style.get(k) == style_00.get(k)


def test_invalid_transfer_mode_raises(table_pair):
    """Test that invalid transfer mode raises ValueError."""
    src_shape, dst_shape = table_pair
    with pytest.raises(ValueError):
        transfer_table_style(src_shape, dst_shape, mode="invalid")