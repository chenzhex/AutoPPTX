import pytest
from pptx import Presentation

from autopptx.Image.style import (
    set_image_style,
    extract_image_style,
    transfer_image_style,
)
from autopptx.Type.find import find_placeholders


TEST_PPTX = "data/output_demo.pptx"


@pytest.fixture
def prs():
    """Load and return the test PowerPoint presentation."""
    return Presentation(TEST_PPTX)


@pytest.fixture
def image_shapes(prs):
    """
    Get two image shapes from slide[3] for testing.

    Returns:
        tuple: (source_image, target_image)
    """
    slide = prs.slides[3]
    images = find_placeholders(slide, "image")
    if len(images) < 2:
        pytest.skip("⚠️ Not enough image placeholders in slide 3.")
    return images[0], images[1]


def test_set_image_style_and_extract(image_shapes):
    """
    Test setting image position, size, rotation, and extracting the result.
    """
    image, _ = image_shapes
    set_image_style(
        image,
        left=2.0,
        top=1.5,
        width=3.0,
        height=2.5,
        rotation=45.0,
    )
    style = extract_image_style(image)
    assert style["left"] == 2.0
    assert style["top"] == 1.5
    assert style["width"] == 3.0
    assert style["height"] == 2.5
    assert style["rotation"] == 45.0


def test_transfer_image_style(image_shapes):
    """
    Test transferring style from one image to another.
    Verifies that target shape matches source style.
    """
    src, dst = image_shapes
    set_image_style(src, left=1.2, top=0.8, width=2.4, height=1.8, rotation=15.0)
    transfer_image_style(src, dst)
    src_style = extract_image_style(src)
    dst_style = extract_image_style(dst)

    assert src_style == dst_style


def test_extract_image_style_invalid_input():
    """
    Test that extract_image_style raises ValueError for invalid input.
    """
    with pytest.raises(ValueError):
        extract_image_style("not a shape")


def test_set_image_style_invalid_input():
    """
    Test that set_image_style raises ValueError for invalid input.
    """
    with pytest.raises(ValueError):
        set_image_style("not an image")


def test_transfer_image_style_invalid_inputs():
    """
    Test that transfer_image_style raises ValueError if either shape is invalid.
    """
    with pytest.raises(ValueError):
        transfer_image_style("not an image", "still not an image")