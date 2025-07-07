import pytest
from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER

from autopptx.Table.table import (
    replace_table,
    replace_table_cell,
    validate_single_table,
)
from autopptx.Table.tables import replace_tables
from autopptx.Type.find import find_placeholders

TEST_PPTX = "data/output_demo.pptx"


@pytest.fixture
def test_slide():
    """Return slide[5] from test presentation."""
    prs = Presentation(TEST_PPTX)
    return prs.slides[5]


@pytest.fixture
def table_data():
    """Return a valid sample table."""
    return [
        ["Header1", "Header2", "Header3"],
        ["Row1", "Value1", "UserA"],
        ["Row2", "Value2", "UserB"],
    ]


@pytest.fixture
def multi_table_data():
    """Return a list of valid 2D tables."""
    return [
        [["A", "B"], ["1", "2"]],
        [["X", "Y", "Z"], ["10", "20", "30"]],
    ]


def test_validate_single_table_valid(table_data):
    """Test valid 2D table format."""
    assert validate_single_table(table_data) is True


@pytest.mark.parametrize(
    "bad_table",
    [None, [], [["A", "B"], ["C"]], "Not a table"],
)
def test_validate_single_table_invalid(bad_table):
    """Test various invalid table structures."""
    assert validate_single_table(bad_table) is False


def test_replace_table_valid(test_slide, table_data):
    """Test replacing a table in the first table placeholder."""
    shapes = find_placeholders(test_slide, "table")
    if not shapes:
        pytest.skip("No table placeholder found in test slide.")

    shape = shapes[0]
    replace_table(shape, table_data)

    table = shape.table
    assert table.cell(0, 0).text == "Header1"
    assert table.cell(2, 2).text == "UserB"


def test_replace_table_invalid_shape(test_slide, table_data):
    """Expect ValueError if shape is not table placeholder."""
    all_shapes = list(test_slide.placeholders)
    non_table_shapes = [
        s for s in all_shapes if s.placeholder_format.type != PP_PLACEHOLDER.TABLE
    ]
    if not non_table_shapes:
        pytest.skip("No non-table placeholder available.")

    with pytest.raises(ValueError):
        replace_table(non_table_shapes[0], table_data)


def test_replace_table_cell_success(test_slide, table_data):
    """Test updating a specific table cell."""
    shapes = find_placeholders(test_slide, "table")
    if not shapes:
        pytest.skip("No table placeholder found.")

    shape = shapes[0]
    replace_table(shape, table_data)
    replace_table_cell(shape, 1, 2, "✅ Passed")

    assert shape.table.cell(1, 2).text == "✅ Passed"


def test_replace_table_cell_out_of_bounds(test_slide, table_data):
    """Test that out-of-range cell raises IndexError."""
    shapes = find_placeholders(test_slide, "table")
    if not shapes:
        pytest.skip("No table placeholder found.")

    shape = shapes[0]
    replace_table(shape, table_data)

    with pytest.raises(IndexError):
        replace_table_cell(shape, 10, 10, "Invalid")


def test_replace_tables_batch(test_slide, multi_table_data):
    """Test replacing multiple table placeholders at once."""
    shapes = find_placeholders(test_slide, "table")
    if not shapes:
        pytest.skip("No table placeholder found.")

    replace_tables(test_slide, multi_table_data)

    assert shapes[0].table.cell(0, 0).text == "A"
    if len(shapes) > 1 and len(multi_table_data) > 1:
        assert shapes[1].table.cell(1, 2).text == "30"


def test_replace_tables_invalid_data(test_slide):
    """Test replace_tables raises error with invalid input."""
    with pytest.raises(ValueError):
        replace_tables(test_slide, "not a list")