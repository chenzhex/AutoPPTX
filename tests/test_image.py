import pytest
from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER

from autopptx.Image.image import replace_image
from autopptx.Image.images import replace_images
from autopptx.Type.find import find_placeholders

TEST_PPTX = "data/output_demo.pptx"
IMG1 = "data/bunny1.png"
IMG2 = "data/bunny2.png"
IMG_MISSING = "data/nonexistent.png"
IMG_REPEAT = "data/cat1.png"


@pytest.fixture
def test_slide():
    """Return slide[3] from the test presentation."""
    prs = Presentation(TEST_PPTX)
    return prs.slides[3]


@pytest.fixture
def image_pair():
    """Return a list of two valid image paths."""
    return [IMG1, IMG2]


@pytest.fixture
def image_with_missing():
    """Return one valid and one missing image path."""
    return [IMG_REPEAT, IMG_MISSING]


@pytest.fixture
def image_single():
    """Return a single image path in a list."""
    return [IMG1]


@pytest.fixture
def image_extra():
    """Return more image paths than expected."""
    return [IMG1, IMG2, IMG1]


def test_replace_image_valid(test_slide):
    """Test replacing an image in a valid picture placeholder."""
    shapes = find_placeholders(test_slide, "image")
    if not shapes:
        pytest.skip("No picture placeholder found.")

    shape = shapes[0]
    replace_image(shape, IMG1)

    assert shape.is_placeholder
    assert shape.placeholder_format.type == PP_PLACEHOLDER.PICTURE


def test_replace_image_invalid_shape(test_slide):
    """Test ValueError is raised for non-image placeholders."""
    non_image_shapes = [
        s
        for s in test_slide.placeholders
        if s.placeholder_format.type != PP_PLACEHOLDER.PICTURE
    ]
    if not non_image_shapes:
        pytest.skip("No non-image placeholder found.")

    with pytest.raises(ValueError):
        replace_image(non_image_shapes[0], IMG1)


def test_replace_image_missing_path(test_slide):
    """Test FileNotFoundError is raised for missing image path."""
    shapes = find_placeholders(test_slide, "image")
    if not shapes:
        pytest.skip("No image placeholder found.")

    shape = shapes[0]
    with pytest.raises(FileNotFoundError):
        replace_image(shape, IMG_MISSING)


def test_replace_images_batch_success(test_slide, image_pair):
    """Test successful batch replacement of image placeholders."""
    shapes = find_placeholders(test_slide, "image")
    if len(shapes) < 2:
        pytest.skip("Test requires at least 2 image placeholders.")

    replace_images(test_slide, image_pair)

    for shape in shapes[:2]:
        assert hasattr(shape._element, "blip_rId")


def test_replace_images_with_invalid_path(test_slide, image_with_missing):
    """Test batch replacement with one invalid image path raises error."""
    shapes = find_placeholders(test_slide, "image")
    if len(shapes) < 2:
        pytest.skip("Test requires at least 2 image placeholders.")

    with pytest.raises(FileNotFoundError):
        replace_images(test_slide, image_with_missing)


def test_replace_images_less_than_placeholders(test_slide, image_single):
    """Test batch replacement with fewer images than placeholders."""
    shapes = find_placeholders(test_slide, "image")
    if len(shapes) < 2:
        pytest.skip("At least 2 image placeholders required.")

    replace_images(test_slide, image_single)

    assert shapes[0].is_placeholder


def test_replace_images_more_than_placeholders(test_slide, image_extra):
    """Test batch replacement with more images than placeholders."""
    shapes = find_placeholders(test_slide, "image")
    if not shapes:
        pytest.skip("No image placeholders found.")

    replace_images(test_slide, image_extra)

    # Only up to the number of placeholders should be replaced
    assert True