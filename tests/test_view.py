import os
import pytest
from pptx import Presentation

from autopptx.View.view import get_table, get_text, get_image, view_slide

TEST_PPTX = "data/output_demo.pptx"
OUTPUT_DIR = "data/extracted_images"


@pytest.fixture(scope="module")
def presentation():
    return Presentation(TEST_PPTX)


def get_first_shape(slide, condition):
    """
    Return the first shape in a slide that satisfies the condition.
    """
    for shape in slide.shapes:
        if condition(shape):
            return shape
    return None


def find_slide(presentation, condition):
    """
    Return the first slide that contains at least one shape satisfying the condition.
    """
    for slide in presentation.slides:
        if any(condition(shape) for shape in slide.shapes):
            return slide
    return None


def test_get_text(presentation):
    """
    Test extracting text from a shape with a text frame.
    """
    slide = find_slide(presentation, lambda s: s.has_text_frame)
    if not slide:
        pytest.skip("No slide with text shape found.")

    shape = get_first_shape(slide, lambda s: s.has_text_frame)
    text = get_text(shape)

    assert isinstance(text, str)
    assert len(text) >= 0


def test_get_table(presentation):
    """
    Test extracting table data from a table shape.
    """
    slide = find_slide(presentation, lambda s: hasattr(s, "table"))
    if not slide:
        pytest.skip("No slide with table shape found.")

    shape = get_first_shape(slide, lambda s: hasattr(s, "table"))
    table = get_table(shape)

    assert isinstance(table, list)
    assert all(isinstance(row, list) for row in table)


def test_get_image(presentation):
    """
    Test extracting and saving an image from a shape.
    """
    slide = find_slide(presentation, lambda s: hasattr(s, "image"))
    if not slide:
        pytest.skip("No slide with image shape found.")

    shape = get_first_shape(slide, lambda s: hasattr(s, "image"))
    path = get_image(shape, output_dir=OUTPUT_DIR)

    assert os.path.exists(path)
    assert path.lower().endswith((".png", ".jpg", ".jpeg", ".bmp", ".gif"))

    os.remove(path)  # Clean up temp file


def test_view_slide(presentation):
    """
    Test extracting structured content from the first slide.
    """
    slide = presentation.slides[0]
    results = view_slide(slide)

    assert isinstance(results, list)
    for item in results:
        assert isinstance(item, dict)
        assert "index" in item
        assert "type" in item
        assert "content" in item