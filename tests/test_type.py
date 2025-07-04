import pytest
from pptx import Presentation

from autopptx.Type.find import find_placeholders
from autopptx.Type.type import (
    is_text,
    is_title,
    is_subtitle,
    is_bodytext,
    is_image,
    is_table,
)

TEST_PPTX = "data/output_demo.pptx"


@pytest.fixture
def test_slide():
    """Return slide[4] from the test presentation."""
    prs = Presentation(TEST_PPTX)
    return prs.slides[4]


@pytest.fixture
def test_shapes(test_slide):
    """Return all shapes from the test slide."""
    return test_slide.shapes


def test_shape_type_classification(test_shapes):
    """Test all shape type utilities return bool on each shape."""
    for shape in test_shapes:
        assert isinstance(is_text(shape), bool)
        assert isinstance(is_title(shape), bool)
        assert isinstance(is_subtitle(shape), bool)
        assert isinstance(is_bodytext(shape), bool)
        assert isinstance(is_image(shape), bool)
        assert isinstance(is_table(shape), bool)


@pytest.mark.parametrize(
    "ptype, checker",
    [
        ("title", is_title),
        ("subtitle", is_subtitle),
        ("bodytext", is_bodytext),
        ("image", is_image),
        ("table", is_table),
    ],
)
def test_placeholder_detection_by_type(test_slide, ptype, checker):
    """Ensure each is_xxx function correctly identifies its placeholder type."""
    shapes = find_placeholders(test_slide, ptype)
    assert isinstance(shapes, list)
    for shape in shapes:
        assert shape.is_placeholder
        assert checker(shape) is True


def test_find_placeholders_invalid(test_slide):
    """Test invalid placeholder type raises ValueError."""
    with pytest.raises(ValueError):
        find_placeholders(test_slide, "unknown_type")