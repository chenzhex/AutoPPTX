import pytest
from pptx import Presentation

from autopptx.Text.text import replace_bodytext
from autopptx.Text.texts import (
    replace_title,
    replace_subtitle,
    replace_bodytexts,
)
from autopptx.Type.find import find_placeholders

TEST_PPTX = "data/output_demo.pptx"


@pytest.fixture
def test_slide():
    """Return slide[0] from test presentation."""
    prs = Presentation(TEST_PPTX)
    return prs.slides[0]


@pytest.mark.parametrize(
    "input_text", ["Updated Title", "Marketing Report"]
)
def test_replace_title(test_slide, input_text):
    """Test replacing title text."""
    replace_title(test_slide, input_text)
    shapes = find_placeholders(test_slide, "title")
    assert shapes and shapes[0].text == input_text


@pytest.mark.parametrize(
    "input_text", ["Subtitle A", "2025 Q1 Summary"]
)
def test_replace_subtitle(test_slide, input_text):
    """Test replacing subtitle text."""
    replace_subtitle(test_slide, input_text)
    shapes = find_placeholders(test_slide, "subtitle")
    assert shapes and shapes[0].text == input_text


@pytest.mark.parametrize(
    "input_text", ["Single body content", "Overview section"]
)
def test_replace_bodytext_valid(test_slide, input_text):
    """Test replacing a single body placeholder."""
    shapes = find_placeholders(test_slide, "bodytext")
    if shapes:
        replace_bodytext(shapes[0], input_text)
        assert shapes[0].text == input_text
    else:
        pytest.skip("No bodytext placeholder found in slide")


def test_replace_bodytext_invalid_shape(test_slide):
    """Test replace_bodytext raises error for non-body placeholder."""
    title_shapes = find_placeholders(test_slide, "title")
    if title_shapes:
        with pytest.raises(ValueError):
            replace_bodytext(title_shapes[0], "Invalid")
    else:
        pytest.skip("No title placeholder found for negative test")


def test_replace_bodytexts_single_box(test_slide):
    """Test inserting multiple paragraphs into one body placeholder."""
    paragraphs = ["Line 1", "Line 2"]
    replace_bodytexts(test_slide, paragraphs, distribute_to_multiple_boxes=False)

    shapes = find_placeholders(test_slide, "bodytext")
    if shapes:
        tf = shapes[0].text_frame
        texts = [p.text for p in tf.paragraphs]
        assert texts == paragraphs
    else:
        pytest.skip("No bodytext placeholder found in slide")


def test_replace_bodytexts_distributed(test_slide):
    """Test distributing paragraphs to multiple text boxes."""
    paragraphs = ["Box A", "Box B", "Box C"]
    replace_bodytexts(test_slide, paragraphs, distribute_to_multiple_boxes=True)

    shapes = find_placeholders(test_slide, "bodytext")
    count = min(len(paragraphs), len(shapes))
    for i in range(count):
        assert shapes[i].text == paragraphs[i]